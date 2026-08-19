from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader


def extract_document(document):
    suffix = Path(document.original_filename).suffix.lower()
    with document.file.open("rb") as stream:
        if suffix == ".pdf":
            reader = PdfReader(stream)
            text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
            return text.strip(), len(reader.pages)
        if suffix == ".docx":
            doc = DocxDocument(stream)
            blocks = [paragraph.text for paragraph in doc.paragraphs]
            for table in doc.tables:
                blocks.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
            return "\n".join(value for value in blocks if value.strip()).strip(), 0
        if suffix == ".pptx":
            presentation = Presentation(stream)
            blocks = []
            for number, slide in enumerate(presentation.slides, start=1):
                slide_text = [
                    shape.text for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if slide_text:
                    blocks.append(f"Slide {number}\n" + "\n".join(slide_text))
            return "\n\n".join(blocks).strip(), len(presentation.slides)
        if suffix in {".txt", ".md", ".markdown"}:
            return stream.read().decode("utf-8").strip(), 0
    return "", 0
